import streamlit as st
from langchain_community.document_loaders import PyPDFLoader
import os
import tempfile
import uuid
from pptx import Presentation
import docx

class DocumentProcessor:
    """
    This class handles the processing of various document types, including PDFs, DOCX, PPT, 
    and plain text. The processed content is stored as a list of pages/chunks for further use.
    """

    def __init__(self):
        # Initialize an empty list to store pages from all documents
        self.pages = []

    def ingest_documents(self):
        """
        Orchestrates the document ingestion process by calling different handler methods based on the 
        type of input. After processing, it displays the total number of pages/chunks processed.
        """
        with st.form("Document Input"):
            # Call the handler for PDF uploads
            self.handle_pdf_upload()
            
            # Call the handler for DOCX uploads
            self.handle_docx_upload()

            # Call the handler for PPT uploads
            self.handle_ppt_upload()

            # Call the handler for plain text input
            self.handle_text_input()

            # Add a submit button to process the inputs
            submitted = st.form_submit_button("Submit")

            if submitted:
                st.write(f"Total pages processed: {len(self.pages)}")

    def handle_pdf_upload(self):
        """
        Handles the upload and processing of PDF files. Each PDF file is saved temporarily, processed 
        to extract its pages, and then the temporary file is deleted. The extracted pages are added 
        to the `pages` list.
        """
        uploaded_files = st.file_uploader(
            "Upload your PDFs",
            type=['pdf'],
            accept_multiple_files=True
        )
        
        if uploaded_files:
            for uploaded_file in uploaded_files:
                try:
                    unique_id = uuid.uuid4().hex
                    original_name, file_extension = os.path.splitext(uploaded_file.name)
                    temp_file_name = f"{original_name}_{unique_id}{file_extension}"
                    temp_file_path = os.path.join(tempfile.gettempdir(), temp_file_name)

                    # Write the uploaded PDF to a temporary file
                    with open(temp_file_path, 'wb') as f:
                        f.write(uploaded_file.getvalue())

                    # Load and process the temporary PDF file using PyPDFLoader
                    loader = PyPDFLoader(temp_file_path)
                    document_pages = loader.load()

                    # Add the extracted pages to the `pages` list
                    self.pages.extend(document_pages)
                except Exception as e:
                    st.error(f"Error processing file {uploaded_file.name}: {e}")
                finally:
                    if os.path.exists(temp_file_path):
                        os.unlink(temp_file_path)

    def handle_docx_upload(self):
        """
        Handles the upload and processing of DOCX files. The text from each DOCX file is extracted 
        and split into lines, which are then added to the `pages` list.
        """
        uploaded_files = st.file_uploader(
            "Upload your DOCX files",
            type=['docx'],
            accept_multiple_files=True
        )
        
        if uploaded_files:
            for uploaded_file in uploaded_files:
                try:
                    # Load the DOCX file using python-docx library
                    doc = docx.Document(uploaded_file)

                    # Append each paragraph's text to the `pages` list
                    for para in doc.paragraphs:
                        self.pages.append(para.text)
                except Exception as e:
                    st.error(f"Error processing DOCX file {uploaded_file.name}: {e}")

    def handle_ppt_upload(self):
        """
        Handles the upload and processing of PPT files. The slides are processed to extract text, 
        which is split into lines and added to the `pages` list.
        """
        uploaded_file = st.file_uploader("Upload your PPT files", type=['pptx'])
        
        if uploaded_file:
            try:
                unique_id = uuid.uuid4().hex
                original_name, file_extension = os.path.splitext(uploaded_file.name)
                temp_file_name = f"{original_name}_{unique_id}{file_extension}"
                temp_file_path = os.path.join(tempfile.gettempdir(), temp_file_name)

                # Write the uploaded PPT file to a temporary file
                with open(temp_file_path, 'wb') as f:
                    f.write(uploaded_file.getvalue())

                # Load the presentation and extract text from each slide
                prs = Presentation(temp_file_path)
                for slide in prs.slides:
                    slide_text = " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
                    self.pages.append(slide_text.splitlines())
            except Exception as e:
                st.error(f"Error processing PPT file {uploaded_file.name}: {e}")
            finally:
                if os.path.exists(temp_file_path):
                    os.unlink(temp_file_path)

    def handle_text_input(self):
        """
        Handles the input of plain text. The text is split into lines and added to the `pages` list.
        """
        text_input = st.text_area("Enter your text here")
        
        if text_input:
            self.pages.append(text_input.splitlines())

# Run the Streamlit app if this script is executed directly
if __name__ == "__main__":
    processor = DocumentProcessor()
    processor.ingest_documents()
